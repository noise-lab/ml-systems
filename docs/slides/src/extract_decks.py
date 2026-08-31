#!/usr/bin/env python3
"""
Extract every lecture .pptx in docs/slides/ into per-deck folders:
  NN-Name/_source-extract.md   (per-slide title + body text + speaker notes)
  NN-Name/images/              (every embedded image, named sNN-iMM.ext)

Run from docs/slides/ :  python3 src/extract_decks.py

Deck -> slug mapping follows issue #10. Two special cases:
  - 11-LogisticRegression.pptx is split into 11-LogisticRegression (slides 1-10)
    and 11-SVM (title slide + slides 11..end).
  - 18-State-Space is already authored in Marp (no .pptx) and is handled
    separately (Marp -> Quarto conversion), so it is not extracted here.
  - 11-LogisticRegression-SVM.pptx is a stale duplicate and is skipped.
"""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# (pptx filename, slug, slide-selection)
#   None        -> all slides
#   (a, b)      -> slides a..b inclusive (1-based)
#   "svm"       -> slide 1 (title) + slides 11..end
DECKS = [
    ("01-intro.pptx", "01-Intro", None),
    ("02-security.pptx", "02-Security", None),
    ("03-performance.pptx", "03-Performance", None),
    ("04-resource.pptx", "04-Resource", None),
    ("05-Data.pptx", "05-Data", None),
    ("06-Features.pptx", "06-Features", None),
    ("07-Preparation.pptx", "07-Preparation", None),
    ("08-Pipelines.pptx", "08-Pipelines", None),
    ("09-Probabilistic.pptx", "09-Probabilistic", None),
    ("10-LinearRegression.pptx", "10-LinearRegression", None),
    ("11-LogisticRegression.pptx", "11-LogisticRegression", (1, 10)),
    ("11-LogisticRegression.pptx", "11-SVM", "svm"),
    ("12-Trees-Ensembles.pptx", "12-Trees-Ensembles", None),
    ("13-Deep-Learning.pptx", "13-Deep-Learning", None),
    ("14-nPrint.pptx", "14-nPrint", None),
    ("15-Dimensionality-Reduction.pptx", "15-Dimensionality-Reduction", None),
    ("16-Clustering.pptx", "16-Clustering", None),
    ("17-Diffusion.pptx", "17-Diffusion", None),
    ("19-Generative.pptx", "19-Generative", None),
    ("20-Timeseries.pptx", "20-Timeseries", None),
]


def selected(n_slides, sel):
    idx = list(range(1, n_slides + 1))
    if sel is None:
        return idx
    if sel == "svm":
        return [1] + [i for i in idx if i >= 11]
    a, b = sel
    return [i for i in idx if a <= i <= b]


def extract(pptx, slug, sel):
    prs = Presentation(pptx)
    slides = list(prs.slides)
    keep = set(selected(len(slides), sel))
    os.makedirs(f"{slug}/images", exist_ok=True)
    lines = [f"# {slug} — source extract", f"_From {pptx}; {len(keep)} slides._", ""]
    n_img = 0
    for i, slide in enumerate(slides, 1):
        if i not in keep:
            continue
        lines.append(f"\n### SLIDE {i}")
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        if title:
            lines.append(f"**Title:** {title}")
        for shape in slide.shapes:
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                for para in shape.text.strip().split("\n"):
                    if para.strip():
                        lines.append(f"- {para.strip()}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    n_img += 1
                    fn = f"{slug}/images/s{i:02d}-i{n_img:02d}.{img.ext}"
                    with open(fn, "wb") as f:
                        f.write(img.blob)
                except Exception as e:
                    lines.append(f"  _[image extract failed: {e}]_")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"\n**Speaker notes:** {notes}")
    with open(f"{slug}/_source-extract.md", "w") as f:
        f.write("\n".join(lines) + "\n")
    return len(keep), n_img


if __name__ == "__main__":
    print(f"{'slug':32} slides  images")
    for pptx, slug, sel in DECKS:
        if not os.path.exists(pptx):
            print(f"{slug:32} MISSING {pptx}")
            continue
        ns, ni = extract(pptx, slug, sel)
        print(f"{slug:32} {ns:6}  {ni:6}")
    print("\nDone. 18-State-Space handled separately (Marp->Quarto).")
